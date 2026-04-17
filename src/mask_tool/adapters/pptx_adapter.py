"""PowerPoint文档(.pptx)适配器"""

from pathlib import Path

from pptx import Presentation

from mask_tool.adapters.base import FileAdapter
from mask_tool.models.detection import DetectionResult, DetectionStatus, Location


class PptxAdapter(FileAdapter):
    """PowerPoint文档脱敏适配器"""

    def supported_extensions(self) -> list[str]:
        return [".pptx"]

    def process(self, input_path: Path, output_dir: Path) -> Path:
        """
        处理PowerPoint文档：
        - 遍历每张幻灯片的所有形状
        - 处理文本框、表格、形状中的文本
        """
        prs = Presentation(str(input_path))
        file_name = input_path.stem
        output_path = output_dir / f"{file_name}_masked.pptx"

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self._process_text_frame(
                        shape.text_frame, str(input_path),
                        slide_index=slide_idx,
                    )
                if shape.has_table:
                    self._process_table(
                        shape.table, str(input_path),
                        slide_index=slide_idx,
                    )

        output_dir.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path

    def _process_text_frame(self, text_frame, file_path: str, slide_index: int) -> None:
        """处理文本框中的所有段落"""
        for para in text_frame.paragraphs:
            full_text = para.text
            if not full_text.strip():
                continue

            # 检测
            results = self.detector.detect(full_text, file_path)
            for r in results:
                r.location.slide = slide_index

            # 策略决策
            results = self.policy.apply(results)

            # 在run级别替换
            for result in results:
                if result.status not in (DetectionStatus.AUTO_MASK, DetectionStatus.SUGGEST_MASK):
                    continue
                if result.text not in full_text:
                    continue

                if self.masker.irreversible:
                    replacement = "***"
                else:
                    from mask_tool.models.mapping import TokenMapping
                    replacement = self.masker.token_gen.generate(result.text, result.text_type)
                    self.masker.mappings.append(TokenMapping(
                        token=replacement,
                        original=result.text,
                        text_type=result.text_type,
                        confidence=result.confidence,
                    ))

                for run in para.runs:
                    if result.text in run.text:
                        run.text = run.text.replace(result.text, replacement)

    def _process_table(self, table, file_path: str, slide_index: int) -> None:
        """处理表格"""
        for row_idx, row in enumerate(table.rows):
            for cell_idx, cell in enumerate(row.cells):
                if cell.text_frame:
                    self._process_text_frame(
                        cell.text_frame, file_path,
                        slide_index=slide_index,
                    )
