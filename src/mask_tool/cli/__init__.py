"""CLI命令行入口"""

import json
import shutil
import time
from pathlib import Path
from typing import List, Optional

import typer
from rich.console import Console
from rich.table import Table

from mask_tool import __version__
from mask_tool.models.config import MaskConfig
from mask_tool.models.detection import DetectionStatus
from mask_tool.core.pipeline import Pipeline

app = typer.Typer(
    name="mask-tool",
    help="本地文件脱敏工具 - 支持可逆脱敏、智能识别、多格式办公文件处理",
    no_args_is_help=True,
)
console = Console()


def _load_config(config_path: Optional[Path], mode: str) -> MaskConfig:
    """加载配置"""
    if config_path and config_path.exists():
        cfg = MaskConfig.from_yaml(config_path)
        if mode != "smart":
            cfg.mode = mode
        return cfg
    # 尝试加载默认配置
    default_paths = [
        Path("config/default.yaml"),
        Path(__file__).parent.parent.parent / "config" / "default.yaml",
    ]
    for p in default_paths:
        if p.exists():
            cfg = MaskConfig.from_yaml(p)
            if mode != "smart":
                cfg.mode = mode
            return cfg
    # 无配置文件时使用默认值
    return MaskConfig(mode=mode)


def _collect_files(input_path: Path) -> List[Path]:
    """收集要处理的文件列表"""
    supported = {".docx", ".xlsx", ".pptx", ".pdf"}
    files = []
    if input_path.is_file():
        if input_path.suffix.lower() in supported:
            files.append(input_path)
        else:
            console.print(f"[yellow]警告: 不支持的文件格式 {input_path.suffix}，已跳过[/yellow]")
    elif input_path.is_dir():
        for f in sorted(input_path.iterdir()):
            if f.is_file() and f.suffix.lower() in supported:
                files.append(f)
    return files


def _print_detection_table(results: list) -> None:
    """用rich表格展示检测结果"""
    table = Table(title="检测结果", show_lines=True)
    table.add_column("敏感信息", style="red bold")
    table.add_column("类别")
    table.add_column("来源")
    table.add_column("置信度", justify="right")
    table.add_column("处置")

    status_map = {
        DetectionStatus.AUTO_MASK: "[green]自动脱敏[/green]",
        DetectionStatus.SUGGEST_MASK: "[yellow]建议脱敏[/yellow]",
        DetectionStatus.HINT_ONLY: "[dim]仅提示[/dim]",
    }

    for r in results:
        table.add_row(
            r.text,
            r.text_type.value,
            r.source,
            f"{r.confidence:.2f}",
            status_map.get(r.status, str(r.status)),
        )

    console.print(table)


@app.command()
def mask(
    input_path: List[Path] = typer.Argument(..., help="输入文件或目录路径", exists=True),
    output: Path = typer.Option("./output", "--output", "-o", help="输出目录"),
    mode: str = typer.Option("smart", "--mode", "-m", help="运行模式: focused/strict/smart/aggressive"),
    config: Optional[Path] = typer.Option(None, "--config", "-c", help="配置文件路径"),
    irreversible: bool = typer.Option(False, "--irreversible", help="使用不可逆脱敏"),
    confirm: bool = typer.Option(False, "--confirm", help="启用交互式确认模式"),
    learn: bool = typer.Option(True, "--learn/--no-learn", help="确认时学习到词库（默认开启）"),
) -> None:
    """对文件执行脱敏处理"""
    console.print(f"[bold green]mask-tool[/bold green] v{__version__}")
    console.print(f"模式: {mode} | 不可逆: {irreversible} | 确认: {'开' if confirm else '关'}")

    cfg = _load_config(config, mode)
    pipeline = Pipeline(cfg)

    if irreversible:
        pipeline.masker = type(pipeline.masker)(
            pipeline.masker.token_gen, irreversible=True,
        )

    # 收集所有文件
    files = []
    for p in input_path:
        files.extend(_collect_files(p))
    if not files:
        console.print("[yellow]未找到可处理的文件[/yellow]")
        raise typer.Exit(1)

    console.print(f"找到 {len(files)} 个文件待处理\n")

    # 确认引擎
    confirm_engine = None
    if confirm:
        from mask_tool.core.confirm import ConfirmEngine
        confirm_engine = ConfirmEngine()

    start_time = time.time()
    for f in files:
        if confirm_engine:
            # 确认模式：先检测，让用户确认，再脱敏
            console.print(f"[bold]  文件: {f.name}[/bold]")
            text = _extract_text(f)
            all_results = pipeline.detector.detect(text, str(f))
            all_results = pipeline.policy.apply(all_results)

            if not all_results:
                console.print("    [dim]未检测到敏感信息，跳过[/dim]\n")
                continue

            # 用户确认
            confirmed = confirm_engine.confirm_batch(all_results, f.name)
            if not confirmed:
                console.print("    [dim]无确认项，跳过脱敏[/dim]\n")
                continue

            # 只对确认的项执行脱敏
            pipeline.report.input_files.append(str(f))
            mappings_before = len(pipeline.masker.mappings)

            # 直接用mask_text处理纯文本（确认模式暂用简化流程）
            masked_text = pipeline.masker.mask_text(text, confirmed)

            # 写入脱敏文件
            _write_masked_file(f, output, masked_text, pipeline)

            console.print()
        else:
            # 自动模式：原有流程
            console.print(f"  处理: {f.name} ...", end=" ")
            try:
                result_path = pipeline.process_file(f, output)
                console.print("[green]✓[/green]")
            except Exception as e:
                console.print(f"[red]✗ {e}[/red]")

    elapsed = time.time() - start_time
    pipeline.report.processing_time_seconds = elapsed

    # 学习机制：将用户确认的词写入词库
    if confirm_engine and learn and confirm_engine.learned:
        _save_learned_words(confirm_engine.get_learned_words(), cfg)

    # 保存映射表
    mapping_path = output / "mapping.json"
    pipeline.save_mapping(mapping_path)
    console.print(f"映射表: {mapping_path}")

    # 保存报告
    report_path = output / "report.json"
    pipeline.save_report(report_path)
    console.print(f"脱敏报告: {report_path}")

    # 打印摘要
    summary = pipeline.report.summary()
    console.print(f"\n[bold]处理摘要:[/bold]")
    console.print(f"  输入文件: {summary['total_input_files']}")
    console.print(f"  自动脱敏: {summary['auto_masked_count']} 项")
    console.print(f"  建议脱敏: {summary['suggested_count']} 项")
    console.print(f"  仅提示: {summary['hint_count']} 项")
    if confirm_engine and confirm_engine.learned:
        console.print(f"  新学词: {len(confirm_engine.learned)} 个")
    console.print(f"  耗时: {elapsed:.2f}s")


@app.command()
def unmask(
    input_path: List[Path] = typer.Argument(..., help="脱敏后的文件或目录路径", exists=True),
    mapping: Path = typer.Option(..., "--mapping", help="映射表文件路径(JSON)", exists=True),
    output: Path = typer.Option("./restored", "--output", "-o", help="输出目录"),
) -> None:
    """反脱敏：根据映射表还原原文"""
    console.print(f"[bold green]mask-tool[/bold green] v{__version__}")

    # 加载映射表
    with open(mapping, "r", encoding="utf-8") as f:
        data = json.load(f)

    tokens = data.get("tokens", {})
    if not tokens:
        console.print("[yellow]映射表为空[/yellow]")
        raise typer.Exit(1)

    console.print(f"加载了 {len(tokens)} 条映射")

    # 收集所有文件
    supported = {".docx", ".xlsx", ".pptx"}
    files = []
    for p in input_path:
        if p.is_file():
            if p.suffix.lower() in supported:
                files.append(p)
        elif p.is_dir():
            for f in sorted(p.iterdir()):
                if f.is_file() and f.suffix.lower() in supported:
                    files.append(f)

    if not files:
        console.print("[yellow]未找到可处理的文件[/yellow]")
        raise typer.Exit(1)

    console.print(f"找到 {len(files)} 个文件待还原\n")
    output.mkdir(parents=True, exist_ok=True)

    for f in files:
        console.print(f"  还原: {f.name} ...", end=" ")
        try:
            suffix = f.suffix.lower()
            output_path = output / f.name
            if suffix == ".docx":
                _unmask_docx(f, output_path, tokens)
            elif suffix == ".xlsx":
                _unmask_xlsx(f, output_path, tokens)
            elif suffix == ".pptx":
                _unmask_pptx(f, output_path, tokens)
            else:
                console.print("[yellow]跳过[/yellow]")
                continue
            console.print("[green]✓[/green]")
        except Exception as e:
            console.print(f"[red]✗ {e}[/red]")

    console.print(f"\n[green]反脱敏完成，输出到: {output}/[/green]")


def _unmask_docx(input_path: Path, output_path: Path, tokens: dict) -> None:
    """反脱敏Word文档"""
    from docx import Document
    shutil.copy2(input_path, output_path)
    doc = Document(str(output_path))
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            for token, info in tokens.items():
                if token in run.text:
                    run.text = run.text.replace(token, info["original"])
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        for token, info in tokens.items():
                            if token in run.text:
                                run.text = run.text.replace(token, info["original"])
    doc.save(str(output_path))


def _unmask_xlsx(input_path: Path, output_path: Path, tokens: dict) -> None:
    """反脱敏Excel文档"""
    from openpyxl import load_workbook
    shutil.copy2(input_path, output_path)
    wb = load_workbook(str(output_path))
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    for token, info in tokens.items():
                        if token in cell.value:
                            cell.value = cell.value.replace(token, info["original"])
    wb.save(str(output_path))


def _unmask_pptx(input_path: Path, output_path: Path, tokens: dict) -> None:
    """反脱敏PPT文档"""
    from pptx import Presentation
    shutil.copy2(input_path, output_path)
    prs = Presentation(str(output_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for token, info in tokens.items():
                            if token in run.text:
                                run.text = run.text.replace(token, info["original"])
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    for token, info in tokens.items():
                                        if token in run.text:
                                            run.text = run.text.replace(token, info["original"])
    prs.save(str(output_path))


@app.command()
def inspect(
    input_path: List[Path] = typer.Argument(..., help="输入文件或目录路径", exists=True),
    config: Optional[Path] = typer.Option(None, "--config", "-c", help="配置文件路径"),
    mode: str = typer.Option("smart", "--mode", "-m", help="运行模式: focused/strict/smart/aggressive"),
) -> None:
    """检测文件中的敏感信息（不执行脱敏）"""
    console.print(f"[bold green]mask-tool[/bold green] v{__version__}")
    console.print(f"模式: {mode}")

    cfg = _load_config(config, mode)
    pipeline = Pipeline(cfg)

    files = []
    for p in input_path:
        files.extend(_collect_files(p))
    if not files:
        console.print("[yellow]未找到可处理的文件[/yellow]")
        raise typer.Exit(1)

    all_results = []
    for f in files:
        console.print(f"  检测: {f.name} ...", end=" ")
        try:
            # 对纯文本文件直接检测
            text = _extract_text(f)
            results = pipeline.detector.detect(text, str(f))
            results = pipeline.policy.apply(results)
            all_results.extend(results)
            console.print(f"发现 {len(results)} 项")
        except Exception as e:
            console.print(f"[red]✗ {e}[/red]")

    if all_results:
        # 按置信度降序排列
        all_results.sort(key=lambda r: r.confidence, reverse=True)
        _print_detection_table(all_results)

        # 统计
        auto = sum(1 for r in all_results if r.status == DetectionStatus.AUTO_MASK)
        suggest = sum(1 for r in all_results if r.status == DetectionStatus.SUGGEST_MASK)
        hint = sum(1 for r in all_results if r.status == DetectionStatus.HINT_ONLY)
        console.print(f"\n[bold]统计: 自动脱敏 {auto} | 建议脱敏 {suggest} | 仅提示 {hint}[/bold]")
    else:
        console.print("\n[green]未检测到敏感信息[/green]")


def _extract_text(file_path: Path) -> str:
    """从文件中提取纯文本"""
    suffix = file_path.suffix.lower()
    if suffix == ".docx":
        from docx import Document
        doc = Document(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs)
    elif suffix == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(str(file_path), read_only=True)
        texts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        texts.append(cell.value)
        wb.close()
        return "\n".join(texts)
    elif suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(file_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)
    elif suffix == ".pdf":
        import fitz
        doc = fitz.open(str(file_path))
        texts = []
        for page in doc:
            texts.append(page.get_text())
        doc.close()
        return "\n".join(texts)
    return ""


def _write_masked_file(
    input_path: Path,
    output_dir: Path,
    masked_text: str,
    pipeline,
) -> Path:
    """将脱敏后的文本写入文件（确认模式的简化版文件输出）

    注意：调用此函数前，pipeline.masker.mask_text() 必须已被调用过，
    以确保 mappings 中已有 original -> token 的映射关系。
    """
    suffix = input_path.suffix.lower()
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"{input_path.stem}_masked{suffix}"

    # 构建 original -> token 替换映射
    replace_map = {m.original: m.token for m in pipeline.masker.get_mappings()}

    if suffix == ".docx":
        from docx import Document
        doc = Document(str(input_path))
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                for original, token in replace_map.items():
                    if original in run.text:
                        run.text = run.text.replace(original, token)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            for original, token in replace_map.items():
                                if original in run.text:
                                    run.text = run.text.replace(original, token)
        doc.save(str(output_path))
    elif suffix == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(str(input_path))
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        for original, token in replace_map.items():
                            if original in cell.value:
                                cell.value = cell.value.replace(original, token)
        wb.save(str(output_path))
    elif suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(input_path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            for original, token in replace_map.items():
                                if original in run.text:
                                    run.text = run.text.replace(original, token)
        prs.save(str(output_path))
    else:
        # 其他格式：直接写纯文本
        output_path.write_text(masked_text, encoding="utf-8")

    pipeline.report.output_files.append(str(output_path))
    console.print(f"    [green]✓ 已保存: {output_path.name}[/green]")
    return output_path


def _save_learned_words(
    learned: dict,
    config: MaskConfig,
) -> None:
    """将学习到的词追加到词库文件"""
    import yaml
    from pathlib import Path

    lexicon_path = Path(config.lexicon_path)
    if not lexicon_path.exists():
        console.print("  [yellow]词库文件不存在，跳过学习写入[/yellow]")
        return

    # 加载现有词库
    with open(lexicon_path, "r", encoding="utf-8") as f:
        existing = yaml.safe_load(f) or {}

    # 合并新词
    new_count = 0
    for category, words in learned.items():
        if category not in existing:
            existing[category] = []
        for word in words:
            if word not in existing[category]:
                existing[category].append(word)
                new_count += 1

    if new_count > 0:
        with open(lexicon_path, "w", encoding="utf-8") as f:
            yaml.dump(existing, f, allow_unicode=True, default_flow_style=False)
        console.print(f"  [blue]✓ {new_count} 个新词已写入词库: {lexicon_path}[/blue]")
    else:
        console.print("  [dim]所有词已存在于词库中，无需更新[/dim]")


@app.command("config")
def config_init(
    output: Path = typer.Option(".", "--output", "-o", help="输出目录"),
) -> None:
    """生成默认配置文件和示例词库"""
    console.print(f"[bold green]mask-tool[/bold green] v{__version__}")

    import shutil
    source_dir = Path(__file__).parent.parent.parent / "config"
    target_dir = Path(output) / "config"
    target_dir.mkdir(parents=True, exist_ok=True)

    for f in source_dir.iterdir():
        target = target_dir / f.name
        if not target.exists():
            shutil.copy2(f, target)
            console.print(f"  创建: {target}")
        else:
            console.print(f"  跳过(已存在): {target}")

    console.print(f"\n[green]配置文件已生成到 {target_dir}/[/green]")
    console.print("编辑 config/default.yaml 自定义配置")
    console.print("编辑 config/sample_lexicon.yaml 管理词库")


@app.command()
def version() -> None:
    """显示版本信息"""
    console.print(f"mask-tool v{__version__}")
