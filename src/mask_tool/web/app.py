"""Streamlit Web 界面 - mask-tool 可视化脱敏工具"""

import json
import random
import shutil
import string
import tempfile
import time
import zipfile
from dataclasses import dataclass, field, asdict
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Optional, Set, Tuple

import streamlit as st
import pandas as pd
import yaml

from mask_tool import __version__
from mask_tool.models.config import MaskConfig
from mask_tool.models.detection import (
    DetectionResult, DetectionStatus, DetectionType, Location,
)
from mask_tool.core.pipeline import Pipeline

# ──────────────────────────────────────────────
# 常量与映射
# ──────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {".docx", ".xlsx", ".pptx", ".pdf"}

TYPE_LABELS: Dict[DetectionType, str] = {
    DetectionType.COMPANY: "🏢 公司/机构",
    DetectionType.GOVERNMENT: "🏛️ 政府",
    DetectionType.PERSON: "👤 人名",
    DetectionType.PROJECT: "📋 项目",
    DetectionType.SUBJECT: "📌 主题",
    DetectionType.LOCATION: "📍 地名",
    DetectionType.AMOUNT: "💰 金额",
    DetectionType.CUSTOM: "🏷️ 自定义",
}

STATUS_LABELS: Dict[DetectionStatus, str] = {
    DetectionStatus.AUTO_MASK: "✅ 自动脱敏",
    DetectionStatus.SUGGEST_MASK: "⚠️ 建议脱敏",
    DetectionStatus.HINT_ONLY: "ℹ️ 仅提示",
}

SOURCE_LABELS: Dict[str, str] = {
    "dictionary": "📘 词典",
    "ner": "🤖 NER",
    "regex": "🔍 正则",
}

MODE_DESCRIPTIONS = {
    "focused": "精准模式：仅自动脱敏词典匹配项（置信度≥0.95），其余仅提示",
    "strict": "严格模式：高置信度自动脱敏，中置信度建议脱敏",
    "smart": "智能模式：平衡自动脱敏与建议脱敏（推荐）",
    "aggressive": "激进模式：尽可能多地脱敏，适合AI前处理",
}

HISTORY_PATH = Path.home() / ".mask-tool" / "history.json"
MAX_HISTORY = 50


# ──────────────────────────────────────────────
# 页面配置
# ──────────────────────────────────────────────

st.set_page_config(
    page_title="mask-tool 文件脱敏工具",
    page_icon="🔒",
    layout="wide",
    initial_sidebar_state="expanded",
)


# ──────────────────────────────────────────────
# 自定义 CSS
# ──────────────────────────────────────────────

def _inject_css():
    st.markdown("""
    <style>
    /* 整体 */
    .main .block-container { padding-top: 2rem; }
    stApp { background: #f8f9fb; }

    /* 侧边栏 */
    [data-testid="stSidebar"] {
        background: linear-gradient(180deg, #1a1a2e 0%, #16213e 100%);
    }
    [data-testid="stSidebar"] * { color: #e0e0e0 !important; }

    /* 指标卡片 */
    .metric-card {
        background: white;
        border-radius: 12px;
        padding: 1rem 1.5rem;
        box-shadow: 0 1px 3px rgba(0,0,0,0.08);
        text-align: center;
    }
    .metric-card .value {
        font-size: 2rem;
        font-weight: 700;
        color: #1a1a2e;
    }
    .metric-card .label {
        font-size: 0.85rem;
        color: #666;
        margin-top: 0.25rem;
    }

    /* 检测结果表格 */
    .detection-table {
        width: 100%;
        border-collapse: separate;
        border-spacing: 0;
        font-size: 0.9rem;
    }
    .detection-table th {
        background: #1a1a2e;
        color: white;
        padding: 0.6rem 0.8rem;
        text-align: left;
        font-weight: 600;
        position: sticky;
        top: 0;
    }
    .detection-table th:first-child { border-radius: 8px 0 0 0; }
    .detection-table th:last-child { border-radius: 0 8px 0 0; }
    .detection-table td {
        padding: 0.5rem 0.8rem;
        border-bottom: 1px solid #eee;
        vertical-align: middle;
    }
    .detection-table tr:hover td { background: #f0f4ff; }
    .detection-table .text-cell {
        font-weight: 600;
        color: #c0392b;
        max-width: 200px;
        overflow: hidden;
        text-overflow: ellipsis;
        white-space: nowrap;
    }
    .detection-table .context-cell {
        color: #555;
        font-size: 0.82rem;
        max-width: 350px;
        overflow: hidden;
        text-overflow: ellipsis;
        white-space: nowrap;
    }
    .detection-table .confidence-high { color: #27ae60; font-weight: 700; }
    .detection-table .confidence-mid { color: #f39c12; font-weight: 600; }
    .detection-table .confidence-low { color: #95a5a6; }

    /* 文件卡片 */
    .file-card {
        background: white;
        border: 1px solid #e0e0e0;
        border-radius: 10px;
        padding: 0.8rem 1rem;
        margin-bottom: 0.5rem;
        display: flex;
        align-items: center;
        gap: 0.8rem;
        transition: box-shadow 0.2s;
    }
    .file-card:hover { box-shadow: 0 2px 8px rgba(0,0,0,0.1); }
    .file-card .icon { font-size: 1.5rem; }
    .file-card .name { font-weight: 600; color: #333; }
    .file-card .size { color: #999; font-size: 0.82rem; }

    /* 按钮 */
    .stButton > button[kind="primary"] {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        border: none;
        color: white;
        font-weight: 600;
        padding: 0.6rem 1.5rem;
        border-radius: 8px;
        transition: transform 0.1s;
    }
    .stButton > button[kind="primary"]:hover { transform: translateY(-1px); }

    /* 步骤指示器 */
    .step-indicator {
        display: flex;
        align-items: center;
        gap: 0.5rem;
        margin-bottom: 1.5rem;
    }
    .step {
        display: flex;
        align-items: center;
        gap: 0.4rem;
        padding: 0.4rem 0.8rem;
        border-radius: 20px;
        font-size: 0.85rem;
        font-weight: 500;
    }
    .step.active { background: #667eea; color: white; }
    .step.done { background: #27ae60; color: white; }
    .step.pending { background: #eee; color: #999; }
    .step-arrow { color: #ccc; }

    /* 标签页 */
    .tab-container {
        display: flex;
        gap: 0;
        border-bottom: 2px solid #eee;
        margin-bottom: 1rem;
    }
    .tab {
        padding: 0.6rem 1.2rem;
        cursor: pointer;
        font-weight: 500;
        color: #666;
        border-bottom: 2px solid transparent;
        margin-bottom: -2px;
        transition: all 0.2s;
    }
    .tab.active {
        color: #667eea;
        border-bottom-color: #667eea;
    }
    .tab:hover { color: #667eea; }

    /* 成功横幅 */
    .success-banner {
        background: linear-gradient(135deg, #27ae60 0%, #2ecc71 100%);
        color: white;
        padding: 1.5rem 2rem;
        border-radius: 12px;
        text-align: center;
        margin: 1rem 0;
    }
    .success-banner h2 { margin: 0 0 0.5rem 0; }
    .success-banner p { margin: 0; opacity: 0.9; }
    </style>
    """, unsafe_allow_html=True)


# ──────────────────────────────────────────────
# 辅助函数
# ──────────────────────────────────────────────

def _extract_text(file_path: Path) -> str:
    """从文件中提取纯文本"""
    suffix = file_path.suffix.lower()
    if suffix == ".docx":
        from docx import Document
        doc = Document(str(file_path))
        texts = []
        for p in doc.paragraphs:
            texts.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
        return "\n".join(texts)
    elif suffix == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(str(file_path), read_only=True)
        texts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        texts.append(cell.value)
        wb.close()
        return "\n".join(texts)
    elif suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(file_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            texts.append(cell.text)
        return "\n".join(texts)
    elif suffix == ".pdf":
        try:
            import fitz
            doc = fitz.open(str(file_path))
            texts = [page.get_text() for page in doc]
            doc.close()
            return "\n".join(texts)
        except ImportError:
            return ""
    return ""


def _file_icon(suffix: str) -> str:
    """返回文件类型图标"""
    icons = {
        ".docx": "📄", ".xlsx": "📊", ".pptx": "📽️", ".pdf": "📕",
    }
    return icons.get(suffix.lower(), "📁")


def _confidence_class(confidence: float) -> str:
    """返回置信度对应的CSS类名"""
    if confidence >= 0.85:
        return "confidence-high"
    elif confidence >= 0.60:
        return "confidence-mid"
    return "confidence-low"


def _load_config(mode: str, config_path: Optional[str] = None) -> MaskConfig:
    """加载配置，如果用户词库不存在则从示例词库复制"""
    if config_path and Path(config_path).exists():
        cfg = MaskConfig.from_yaml(Path(config_path))
        cfg.mode = mode
        _ensure_lexicon_exists(cfg)
        return cfg

    # 尝试默认配置路径
    default_paths = [
        Path("config/default.yaml"),
        Path(__file__).parent.parent.parent / "config" / "default.yaml",
    ]
    for p in default_paths:
        if p.exists():
            cfg = MaskConfig.from_yaml(p)
            cfg.mode = mode
            _ensure_lexicon_exists(cfg)
            return cfg

    return MaskConfig(mode=mode)


def _ensure_lexicon_exists(cfg: MaskConfig) -> None:
    """如果用户词库文件不存在，从示例词库复制"""
    lexicon_path = Path(cfg.lexicon_path)
    if not lexicon_path.exists():
        sample_path = Path("config/sample_lexicon.yaml")
        if sample_path.exists():
            lexicon_path.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(sample_path, lexicon_path)


def _dedup_results(results: List[DetectionResult]) -> List[DetectionResult]:
    """跨文件去重：按(text, text_type)去重"""
    seen: Set[Tuple[str, str]] = set()
    deduped = []
    for r in results:
        key = (r.text, r.text_type.value)
        if key not in seen:
            seen.add(key)
            deduped.append(r)
    return deduped


def _results_to_dataframe(results: List[DetectionResult]) -> pd.DataFrame:
    """将检测结果转为 DataFrame"""
    rows = []
    for i, r in enumerate(results):
        rows.append({
            "序号": i + 1,
            "敏感信息": r.text,
            "类别": TYPE_LABELS.get(r.text_type, r.text_type.value),
            "类别值": r.text_type.value,
            "来源": SOURCE_LABELS.get(r.source, r.source),
            "置信度": r.confidence,
            "处置": STATUS_LABELS.get(r.status, r.status.value),
            "状态值": r.status.value,
            "文件": Path(r.location.file).name if r.location.file else "",
            "上下文": r.context,
        })
    return pd.DataFrame(rows)


def _do_mask_file(
    input_path: Path,
    output_dir: Path,
    pipeline: Pipeline,
    confirmed_results: List[DetectionResult],
) -> Optional[Path]:
    """对单个文件执行脱敏（基于用户确认的结果）

    流程：先用 mask_text 在纯文本上生成 Token 映射，
    再用 (original -> token) 映射在原始文件上做替换。
    """
    suffix = input_path.suffix.lower()
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"{input_path.stem}_masked{suffix}"

    # 1. 先在纯文本上执行 mask_text，生成 original -> token 映射
    text = _extract_text(input_path)
    if text.strip():
        pipeline.masker.mask_text(text, confirmed_results)

    # 2. 构建替换映射表：original -> token
    replace_map: Dict[str, str] = {}
    for m in pipeline.masker.get_mappings():
        replace_map[m.original] = m.token

    if not replace_map:
        # 没有需要替换的内容，直接复制原文件
        shutil.copy2(input_path, output_path)
        return output_path

    # 3. 在原始文件上执行替换（保持格式）
    if suffix == ".docx":
        from docx import Document
        doc = Document(str(input_path))
        # 段落
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                for original, token in replace_map.items():
                    if original in run.text:
                        run.text = run.text.replace(original, token)
        # 表格
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            for original, token in replace_map.items():
                                if original in run.text:
                                    run.text = run.text.replace(original, token)
        doc.save(str(output_path))

    elif suffix == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(str(input_path))
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        for original, token in replace_map.items():
                            if original in cell.value:
                                cell.value = cell.value.replace(original, token)
        wb.save(str(output_path))

    elif suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(input_path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            for original, token in replace_map.items():
                                if original in run.text:
                                    run.text = run.text.replace(original, token)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        for original, token in replace_map.items():
                                            if original in run.text:
                                                run.text = run.text.replace(original, token)
        prs.save(str(output_path))

    else:
        # PDF 或其他：写纯文本
        masked_text = text
        for original, token in replace_map.items():
            masked_text = masked_text.replace(original, token)
        output_path.write_text(masked_text, encoding="utf-8")

    return output_path


# ──────────────────────────────────────────────
# 反脱敏函数
# ──────────────────────────────────────────────

def _unmask_docx(input_path: Path, output_path: Path, tokens: dict) -> None:
    """反脱敏 docx 文件。tokens 格式: {token_str: original_str}"""
    from docx import Document
    shutil.copy2(input_path, output_path)
    doc = Document(str(output_path))
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            for token, original in tokens.items():
                if token in run.text:
                    run.text = run.text.replace(token, original)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        for token, original in tokens.items():
                            if token in run.text:
                                run.text = run.text.replace(token, original)
    doc.save(str(output_path))


def _unmask_xlsx(input_path: Path, output_path: Path, tokens: dict) -> None:
    """反脱敏 xlsx 文件。tokens 格式: {token_str: original_str}"""
    from openpyxl import load_workbook
    shutil.copy2(input_path, output_path)
    wb = load_workbook(str(output_path))
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    for token, original in tokens.items():
                        if token in cell.value:
                            cell.value = cell.value.replace(token, original)
    wb.save(str(output_path))


def _unmask_pptx(input_path: Path, output_path: Path, tokens: dict) -> None:
    """反脱敏 pptx 文件。tokens 格式: {token_str: original_str}"""
    from pptx import Presentation
    shutil.copy2(input_path, output_path)
    prs = Presentation(str(output_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for token, original in tokens.items():
                            if token in run.text:
                                run.text = run.text.replace(token, original)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    for token, original in tokens.items():
                                        if token in run.text:
                                            run.text = run.text.replace(token, original)
    prs.save(str(output_path))


def _unmask_file(input_path: Path, output_path: Path, tokens: dict) -> Optional[Path]:
    """根据文件类型分发反脱敏。tokens 格式: {token_str: original_str}"""
    suffix = input_path.suffix.lower()
    try:
        if suffix == ".docx":
            _unmask_docx(input_path, output_path, tokens)
        elif suffix == ".xlsx":
            _unmask_xlsx(input_path, output_path, tokens)
        elif suffix == ".pptx":
            _unmask_pptx(input_path, output_path, tokens)
        else:
            # 纯文本文件
            text = input_path.read_text(encoding="utf-8")
            for token, original in tokens.items():
                text = text.replace(token, original)
            output_path.write_text(text, encoding="utf-8")
        return output_path
    except Exception as e:
        st.error(f"反脱敏 {input_path.name} 时出错: {e}")
        return None


# ──────────────────────────────────────────────
# 批次管理
# ──────────────────────────────────────────────

@dataclass
class BatchRecord:
    batch_id: str
    batch_name: str
    created_at: str  # ISO格式
    file_count: int
    mask_count: int
    mapping_data: str  # JSON字符串


def _generate_batch_id() -> str:
    """生成批次ID，格式：MSK-YYYYMMDD-HHMMSS-XXX"""
    now = datetime.now()
    date_part = now.strftime("%Y%m%d")
    time_part = now.strftime("%H%M%S")
    rand_part = "".join(random.choices(string.ascii_uppercase + string.digits, k=3))
    return f"MSK-{date_part}-{time_part}-{rand_part}"


def _load_history() -> List[BatchRecord]:
    """加载批次历史记录"""
    if not HISTORY_PATH.exists():
        return []
    try:
        with open(HISTORY_PATH, "r", encoding="utf-8") as f:
            data = json.load(f)
        return [BatchRecord(**item) for item in data]
    except Exception:
        return []


def _save_history(records: List[BatchRecord]) -> None:
    """保存批次历史记录"""
    HISTORY_PATH.parent.mkdir(parents=True, exist_ok=True)
    # 只保留最近 MAX_HISTORY 条
    records = records[-MAX_HISTORY:]
    data = [asdict(r) for r in records]
    with open(HISTORY_PATH, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def _add_history(record: BatchRecord) -> None:
    """添加一条历史记录"""
    records = _load_history()
    records.append(record)
    _save_history(records)


# ──────────────────────────────────────────────
# 侧边栏
# ──────────────────────────────────────────────

def render_sidebar():
    """渲染侧边栏配置"""
    with st.sidebar:
        st.title("🔒 mask-tool")
        st.caption(f"v{__version__} · 文件脱敏工具")

        st.markdown("---")

        # 运行模式
        st.subheader("⚙️ 运行模式")
        mode = st.selectbox(
            "选择模式",
            options=["focused", "smart", "strict", "aggressive"],
            format_func=lambda x: {
                "focused": "🎯 精准模式",
                "smart": "🧠 智能模式（推荐）",
                "strict": "🔒 严格模式",
                "aggressive": "🚀 激进模式",
            }.get(x, x),
            index=1,
        )
        st.caption(MODE_DESCRIPTIONS.get(mode, ""))

        st.markdown("---")

        # NER 开关
        st.subheader("🤖 NER 引擎")
        ner_enabled = st.toggle(
            "启用 jieba NER",
            value=True,
            help="启用后可识别词典未覆盖的实体（人名、地名、机构名等），但可能产生误识别",
        )

        st.markdown("---")

        # 脱敏选项
        st.subheader("🔧 脱敏选项")
        irreversible = st.checkbox(
            "不可逆脱敏",
            value=False,
            help="启用后将用 *** 替换敏感信息，无法还原",
        )
        learn_words = st.checkbox(
            "学习新词到词库",
            value=True,
            help="确认时标记为'加入词库'的词将写入词库文件",
        )

        st.markdown("---")

        # 词库管理
        st.subheader("📚 词库管理")
        lexicon_info = _get_lexicon_info()
        if lexicon_info:
            st.metric("词库词条", lexicon_info["total"])

            # 2.1: 每个类别可展开查看明细
            lexicon_data = _get_lexicon_data()
            if lexicon_data:
                for cat in sorted(lexicon_data.keys(), key=lambda c: -len(lexicon_data[c])):
                    label = TYPE_LABELS.get(DetectionType(cat), cat)
                    count = len(lexicon_data[cat])
                    with st.expander(f"{label}: {count} 条"):
                        for word in lexicon_data[cat]:
                            st.code(word)
        else:
            st.caption("词库未加载")

        # 2.2: 手动录入词条
        with st.expander("✏️ 手动录入词条", expanded=False):
            valid_categories = {t.value: TYPE_LABELS.get(t, t.value) for t in DetectionType}
            col_cat, col_word = st.columns([1, 2])
            with col_cat:
                input_cat = st.selectbox(
                    "类别",
                    options=list(valid_categories.keys()),
                    format_func=lambda x: valid_categories[x],
                    key="manual_cat",
                    label_visibility="collapsed",
                )
            with col_word:
                input_words = st.text_area(
                    "词条（多条用逗号或换行分隔）",
                    placeholder="输入词条，多条用逗号或换行分隔...",
                    key="manual_words",
                    label_visibility="collapsed",
                    height=70,
                )
            # "其他"类别：允许自定义
            if input_cat == "custom":
                custom_cat_name = st.text_input(
                    "自定义类别名称（留空则归入 custom）",
                    key="custom_cat_name",
                    placeholder="如：brand, department...",
                )
            if st.button("➕ 添加到词库", use_container_width=True, key="add_words_btn"):
                _add_words_to_lexicon(input_cat, input_words, custom_cat_name if input_cat == "custom" else None)

        # 批量导入词库
        with st.expander("📥 批量导入词条", expanded=False):
            st.caption("支持 YAML 或 TXT 格式")
            import_file = st.file_uploader(
                "选择词库文件",
                type=["yaml", "yml", "txt"],
                key="lexicon_upload",
                label_visibility="collapsed",
            )
            if import_file:
                _import_lexicon(import_file)

        st.markdown("---")
        st.caption("mask-tool · MIT License")

    return mode, ner_enabled, irreversible, learn_words


def _get_lexicon_data() -> Optional[Dict[str, List[str]]]:
    """读取词库完整数据（分类别返回词条列表）"""
    try:
        config_paths = [
            Path("config/lexicon.yaml"),
            Path("config/sample_lexicon.yaml"),
            Path(__file__).parent.parent.parent / "config" / "lexicon.yaml",
            Path(__file__).parent.parent.parent / "config" / "sample_lexicon.yaml",
        ]
        for p in config_paths:
            if p.exists():
                with open(p, "r", encoding="utf-8") as f:
                    data = yaml.safe_load(f) or {}
                return {k: v for k, v in data.items() if isinstance(v, list)}
    except Exception:
        pass
    return None


def _add_words_to_lexicon(category: str, words_text: str, custom_category: Optional[str] = None) -> None:
    """手动添加词条到用户词库"""
    if not words_text or not words_text.strip():
        st.warning("请输入词条内容")
        return

    # 确定实际类别
    actual_cat = custom_category.strip() if custom_category and custom_category.strip() else category

    # 解析词条（支持逗号、中文逗号、换行分隔）
    words = []
    for part in words_text.replace("，", ",").replace("\n", ",").split(","):
        w = part.strip()
        if w:
            words.append(w)

    if not words:
        st.warning("未识别到有效词条")
        return

    # 确保用户词库存在
    lexicon_path = Path("config/lexicon.yaml")
    if not lexicon_path.exists():
        sample_path = Path("config/sample_lexicon.yaml")
        if sample_path.exists():
            shutil.copy2(sample_path, lexicon_path)
        else:
            lexicon_path.write_text("", encoding="utf-8")

    # 读取现有词库
    with open(lexicon_path, "r", encoding="utf-8") as f:
        existing = yaml.safe_load(f) or {}

    # 确保类别存在
    if actual_cat not in existing:
        existing[actual_cat] = []

    # 添加新词条（去重）
    added = 0
    for word in words:
        if word not in existing[actual_cat]:
            existing[actual_cat].append(word)
            added += 1

    # 保存
    if added > 0:
        with open(lexicon_path, "w", encoding="utf-8") as f:
            yaml.dump(existing, f, allow_unicode=True, default_flow_style=False)
        st.success(f"✅ 成功添加 {added} 条词条到 [{actual_cat}]")
    else:
        st.info("ℹ️ 所有词条已存在于词库中")


def _get_lexicon_info() -> Optional[dict]:
    """获取词库统计信息（优先读取用户词库 lexicon.yaml）"""
    try:
        # 按优先级查找词库文件
        config_paths = [
            Path("config/lexicon.yaml"),
            Path("config/sample_lexicon.yaml"),
            Path(__file__).parent.parent.parent / "config" / "lexicon.yaml",
            Path(__file__).parent.parent.parent / "config" / "sample_lexicon.yaml",
        ]
        for p in config_paths:
            if p.exists():
                with open(p, "r", encoding="utf-8") as f:
                    data = yaml.safe_load(f) or {}
                categories = {k: len(v) for k, v in data.items() if isinstance(v, list)}
                return {
                    "total": sum(categories.values()),
                    "categories": categories,
                    "path": str(p),
                }
    except Exception:
        pass
    return None


def _import_lexicon(uploaded_file) -> None:
    """从上传的文件批量导入词条到用户词库

    支持格式：
    - YAML: 与 sample_lexicon.yaml 相同格式（{category: [word1, word2, ...]}）
    - TXT: 每行一个词条，格式为 "类别:词条" 或纯词条（默认归入 custom）
    """
    import io

    # 确定用户词库路径
    lexicon_path = Path("config/lexicon.yaml")
    if not lexicon_path.exists():
        # 从示例词库复制
        sample_path = Path("config/sample_lexicon.yaml")
        if sample_path.exists():
            shutil.copy2(sample_path, lexicon_path)
        else:
            lexicon_path.write_text("", encoding="utf-8")

    # 读取现有词库
    with open(lexicon_path, "r", encoding="utf-8") as f:
        existing = yaml.safe_load(f) or {}

    # 确保所有类别键存在
    valid_categories = [t.value for t in DetectionType]
    for cat in valid_categories:
        if cat not in existing:
            existing[cat] = []

    filename = uploaded_file.name.lower()
    added_count = 0

    if filename.endswith((".yaml", ".yml")):
        # YAML 格式导入
        content = uploaded_file.read().decode("utf-8")
        new_data = yaml.safe_load(content)
        if isinstance(new_data, dict):
            for cat, words in new_data.items():
                if isinstance(words, list) and cat in valid_categories:
                    for word in words:
                        if isinstance(word, str) and word.strip() and word not in existing[cat]:
                            existing[cat].append(word.strip())
                            added_count += 1
                elif isinstance(words, list):
                    # 未知类别，归入 custom
                    for word in words:
                        if isinstance(word, str) and word.strip() and word not in existing["custom"]:
                            existing["custom"].append(word.strip())
                            added_count += 1

    elif filename.endswith(".txt"):
        # TXT 格式导入：每行一个词条
        content = uploaded_file.read().decode("utf-8")
        for line in content.strip().splitlines():
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            if ":" in line:
                # 格式：类别:词条
                cat, word = line.split(":", 1)
                cat = cat.strip().lower()
                word = word.strip()
                if cat in valid_categories and word:
                    if word not in existing[cat]:
                        existing[cat].append(word)
                        added_count += 1
            else:
                # 纯词条，归入 custom
                if line not in existing["custom"]:
                    existing["custom"].append(line)
                    added_count += 1

    # 保存
    if added_count > 0:
        with open(lexicon_path, "w", encoding="utf-8") as f:
            yaml.dump(existing, f, allow_unicode=True, default_flow_style=False)
        st.success(f"✅ 成功导入 {added_count} 条新词条到词库")
    else:
        st.info("ℹ️ 没有新词条需要导入（全部已存在）")


# ──────────────────────────────────────────────
# 步骤指示器
# ──────────────────────────────────────────────

def render_steps(current_step: int):
    """渲染步骤指示器 (1-4)"""
    steps = [
        ("1", "📤 上传文件"),
        ("2", "🔍 检测分析"),
        ("3", "✅ 确认选择"),
        ("4", "💾 执行脱敏"),
    ]
    step_html = '<div class="step-indicator">'
    for i, (num, label) in enumerate(steps):
        step_num = i + 1
        cls = "done" if step_num < current_step else ("active" if step_num == current_step else "pending")
        step_html += f'<div class="step {cls}">{label}</div>'
        if i < len(steps) - 1:
            step_html += '<span class="step-arrow">→</span>'
    step_html += '</div>'
    st.markdown(step_html, unsafe_allow_html=True)


# ──────────────────────────────────────────────
# 标签页1：脱敏处理
# ──────────────────────────────────────────────

def _render_masking_tab(mode: str, ner_enabled: bool, irreversible: bool, learn_words: bool):
    """渲染脱敏处理标签页"""

    # 如果已有脱敏结果，展示结果页面
    if "mask_result" in st.session_state:
        _render_mask_result()
        return

    # ── Step 1: 文件上传 ──
    render_steps(1)

    uploaded_files = st.file_uploader(
        "上传待脱敏文件",
        type=["docx", "xlsx", "pptx", "pdf"],
        accept_multiple_files=True,
        label_visibility="collapsed",
    )

    if not uploaded_files:
        st.info("📤 请上传需要脱敏的文件（支持 .docx / .xlsx / .pptx / .pdf）")
        return

    # 显示已上传文件
    st.markdown("#### 📁 已上传文件")
    file_cols = st.columns(min(len(uploaded_files), 4))
    for i, f in enumerate(uploaded_files):
        with file_cols[i % len(file_cols)]:
            icon = _file_icon(Path(f.name).suffix)
            size_kb = f.size / 1024
            st.markdown(
                f'<div class="file-card">'
                f'<span class="icon">{icon}</span>'
                f'<div><div class="name">{f.name}</div>'
                f'<div class="size">{size_kb:.1f} KB</div></div>'
                f'</div>',
                unsafe_allow_html=True,
            )

    # ── Step 2: 检测分析 ──
    st.markdown("---")
    render_steps(2)

    if st.button("🔍 开始检测", type="primary", width="stretch"):
        with st.spinner("正在分析文件，检测敏感信息..."):
            _run_detection(
                uploaded_files, mode, ner_enabled,
            )

    # 检查是否已有检测结果
    if "detection_results" not in st.session_state:
        return

    all_results = st.session_state["detection_results"]
    file_results = st.session_state["file_results"]

    if not all_results:
        st.success("✅ 未检测到敏感信息，文件安全！")
        return

    # ── 检测结果统计 ──
    st.markdown("#### 📊 检测结果概览")

    # 统计卡片
    col1, col2, col3, col4 = st.columns(4)
    auto_count = sum(1 for r in all_results if r.status == DetectionStatus.AUTO_MASK)
    suggest_count = sum(1 for r in all_results if r.status == DetectionStatus.SUGGEST_MASK)
    hint_count = sum(1 for r in all_results if r.status == DetectionStatus.HINT_ONLY)

    with col1:
        st.markdown(
            f'<div class="metric-card"><div class="value">{len(all_results)}</div><div class="label">检测总数</div></div>',
            unsafe_allow_html=True,
        )
    with col2:
        st.markdown(
            f'<div class="metric-card"><div class="value" style="color:#27ae60">{auto_count}</div><div class="label">自动脱敏</div></div>',
            unsafe_allow_html=True,
        )
    with col3:
        st.markdown(
            f'<div class="metric-card"><div class="value" style="color:#f39c12">{suggest_count}</div><div class="label">建议脱敏</div></div>',
            unsafe_allow_html=True,
        )
    with col4:
        st.markdown(
            f'<div class="metric-card"><div class="value" style="color:#95a5a6">{hint_count}</div><div class="label">仅提示</div></div>',
            unsafe_allow_html=True,
        )

    # 类别分布
    type_counts: Dict[str, int] = {}
    for r in all_results:
        label = TYPE_LABELS.get(r.text_type, r.text_type.value)
        type_counts[label] = type_counts.get(label, 0) + 1

    if type_counts:
        chart_cols = st.columns([2, 1])
        with chart_cols[0]:
            # 用原生 HTML 条形图代替 st.bar_chart（避免 pyarrow 依赖）
            sorted_counts = sorted(type_counts.items(), key=lambda x: x[1])
            max_val = max(type_counts.values()) if type_counts else 1
            bars_html = '<div style="font-size:0.85rem;">'
            for label, count in sorted_counts:
                pct = int(count / max_val * 100)
                bars_html += (
                    f'<div style="display:flex;align-items:center;margin-bottom:4px;">'
                    f'<span style="width:120px;flex-shrink:0;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;">{label}</span>'
                    f'<div style="flex:1;background:#eee;border-radius:4px;height:22px;position:relative;">'
                    f'<div style="background:linear-gradient(90deg,#667eea,#764ba2);width:{pct}%;height:100%;border-radius:4px;min-width:2px;"></div>'
                    f'<span style="position:absolute;right:6px;top:2px;font-size:0.78rem;font-weight:600;">{count}</span>'
                    f'</div></div>'
                )
            bars_html += '</div>'
            st.markdown(bars_html, unsafe_allow_html=True)
        with chart_cols[1]:
            st.markdown("**类别分布**")
            for label, count in sorted(type_counts.items(), key=lambda x: -x[1]):
                st.markdown(f"- {label}: **{count}** 项")

    # ── Step 3: 确认选择 ──
    st.markdown("---")
    render_steps(3)

    # 初始化选择状态
    if "user_selections" not in st.session_state:
        # 默认：自动脱敏和建议脱敏的项都勾选
        st.session_state["user_selections"] = {
            i: (r.status in (DetectionStatus.AUTO_MASK, DetectionStatus.SUGGEST_MASK))
            for i, r in enumerate(all_results)
        }

    # 筛选器
    st.markdown("#### 🎛️ 筛选与选择")

    filter_cols = st.columns(5)
    with filter_cols[0]:
        filter_type = st.selectbox(
            "按类别筛选",
            options=["全部"] + list(type_counts.keys()),
            key="filter_type",
        )
    with filter_cols[1]:
        filter_status = st.selectbox(
            "按处置筛选",
            options=["全部", "✅ 自动脱敏", "⚠️ 建议脱敏", "ℹ️ 仅提示"],
            key="filter_status",
        )
    with filter_cols[2]:
        filter_source = st.selectbox(
            "按来源筛选",
            options=["全部", "📘 词典", "🤖 NER", "🔍 正则"],
            key="filter_source",
        )
    with filter_cols[3]:
        filter_file = st.selectbox(
            "按文件筛选",
            options=["全部"] + list(file_results.keys()),
            key="filter_file",
        )
    with filter_cols[4]:
        search_text = st.text_input("搜索", placeholder="输入关键词...", key="search_text")

    # 应用筛选
    filtered_indices = []
    for i, r in enumerate(all_results):
        # 类别筛选
        if filter_type != "全部":
            if TYPE_LABELS.get(r.text_type, r.text_type.value) != filter_type:
                continue
        # 状态筛选
        if filter_status != "全部":
            status_map = {
                "✅ 自动脱敏": DetectionStatus.AUTO_MASK,
                "⚠️ 建议脱敏": DetectionStatus.SUGGEST_MASK,
                "ℹ️ 仅提示": DetectionStatus.HINT_ONLY,
            }
            if r.status != status_map.get(filter_status):
                continue
        # 来源筛选
        if filter_source != "全部":
            source_map = {
                "📘 词典": "dictionary",
                "🤖 NER": "ner",
                "🔍 正则": "regex",
            }
            if r.source != source_map.get(filter_source):
                continue
        # 文件筛选
        if filter_file != "全部":
            if Path(r.location.file).name != filter_file:
                continue
        # 搜索
        if search_text:
            if search_text.lower() not in r.text.lower() and search_text.lower() not in r.context.lower():
                continue
        filtered_indices.append(i)

    # 批量操作按钮
    batch_cols = st.columns(6)
    with batch_cols[0]:
        if st.button("☑️ 全选当前", width="stretch"):
            for i in filtered_indices:
                st.session_state["user_selections"][i] = True
            st.rerun()
    with batch_cols[1]:
        if st.button("☐️ 取消全选", width="stretch"):
            for i in filtered_indices:
                st.session_state["user_selections"][i] = False
            st.rerun()
    with batch_cols[2]:
        if st.button("🔄 反选", width="stretch"):
            for i in filtered_indices:
                st.session_state["user_selections"][i] = not st.session_state["user_selections"][i]
            st.rerun()
    with batch_cols[3]:
        if st.button("✅ 仅选自动脱敏", width="stretch"):
            for i in filtered_indices:
                st.session_state["user_selections"][i] = (
                    all_results[i].status == DetectionStatus.AUTO_MASK
                )
            st.rerun()
    with batch_cols[4]:
        if st.button("⚠️ 仅选建议脱敏", width="stretch"):
            for i in filtered_indices:
                st.session_state["user_selections"][i] = (
                    all_results[i].status == DetectionStatus.SUGGEST_MASK
                )
            st.rerun()
    with batch_cols[5]:
        if st.button("📚 全部加入词库", width="stretch"):
            for i in filtered_indices:
                st.session_state["user_selections"][i] = True
                if "learn_set" not in st.session_state:
                    st.session_state["learn_set"] = set()
                st.session_state["learn_set"].add(i)
            st.rerun()

    # 选中计数
    selected_count = sum(
        1 for i in filtered_indices if st.session_state["user_selections"].get(i, False)
    )
    st.caption(f"当前显示 {len(filtered_indices)} 项，已选中 **{selected_count}** 项")

    # 检测结果表格（使用 AgGrid，勾选不触发 rerun）
    if filtered_indices:
        from st_aggrid import AgGrid, GridOptionsBuilder, GridUpdateMode

        display_rows = []
        for i in filtered_indices:
            r = all_results[i]
            display_rows.append({
                "index": i,
                "选择": st.session_state["user_selections"].get(i, False),
                "敏感信息": r.text,
                "类别": TYPE_LABELS.get(r.text_type, r.text_type.value),
                "来源": SOURCE_LABELS.get(r.source, r.source),
                "置信度": r.confidence,
                "处置": STATUS_LABELS.get(r.status, r.status.value),
                "文件": Path(r.location.file).name if r.location.file else "",
                "上下文": r.context[:80] + "..." if len(r.context) > 80 else r.context,
            })

        df_display = pd.DataFrame(display_rows)

        # 构建 AgGrid 配置
        gb = GridOptionsBuilder.from_dataframe(df_display)
        gb.configure_column("index", hide=True)
        gb.configure_column("选择", header_checkbox_selection=True, editable=True)
        gb.configure_column("敏感信息", editable=False)
        gb.configure_column("类别", editable=False)
        gb.configure_column("来源", editable=False)
        gb.configure_column("置信度", editable=False, type=["numericColumn"], precision_format=2)
        gb.configure_column("处置", editable=False)
        gb.configure_column("文件", editable=False)
        gb.configure_column("上下文", editable=False)
        gb.configure_selection(
            selection_mode="multiple",
            use_checkbox=True,
            pre_selected_rows=[j for j, row in enumerate(display_rows) if row["选择"]],
        )
        gb.configure_pagination(pagination_auto_page_size=False, pagination_page_size=30)
        gridOptions = gb.build()

        # 渲染 AgGrid（fit_column=True 自动列宽）
        grid_response = AgGrid(
            df_display,
            gridOptions=gridOptions,
            update_mode=GridUpdateMode.NO_UPDATE,
            fit_columns_on_grid_load=True,
            height=500,
            allow_unsafe_jscode=True,
            theme="streamlit",
        )

        # 从 AgGrid 响应中同步选择状态
        selected_rows = grid_response.get("selected_rows")
        if selected_rows is not None:
            selected_indices_in_grid = set()
            for row in selected_rows:
                idx = int(row.get("index", -1))
                if idx >= 0:
                    selected_indices_in_grid.add(idx)
            # 更新 session_state
            changed = False
            for i in filtered_indices:
                new_val = i in selected_indices_in_grid
                if st.session_state["user_selections"].get(i) != new_val:
                    st.session_state["user_selections"][i] = new_val
                    changed = True
            # 如果选择状态有变化，更新计数显示（不 rerun 整个页面）
            if changed:
                selected_count = sum(
                    1 for i in filtered_indices if st.session_state["user_selections"].get(i, False)
                )
                st.caption(f"当前显示 {len(filtered_indices)} 项，已选中 **{selected_count}** 项")

    # ── Step 4: 执行脱敏 ──
    st.markdown("---")
    render_steps(4)

    # 最终确认的项
    final_selected = [
        i for i in range(len(all_results))
        if st.session_state["user_selections"].get(i, False)
    ]

    if not final_selected:
        st.warning("⚠️ 请至少选择一项进行脱敏")
        return

    st.markdown(f"#### 📋 即将脱敏 **{len(final_selected)}** 项")

    # 展示选中项预览
    preview_items = []
    for i in final_selected:
        r = all_results[i]
        preview_items.append(f"- {r.text} ({TYPE_LABELS.get(r.text_type, '')})")
    with st.expander("查看选中项详情", expanded=False):
        st.markdown("\n".join(preview_items[:50]))
        if len(preview_items) > 50:
            st.caption(f"... 共 {len(preview_items)} 项")

    # 批次信息
    st.markdown("#### 📦 批次信息")
    batch_cols = st.columns(2)
    with batch_cols[0]:
        batch_name = st.text_input(
            "批次名称（可选）",
            placeholder="例如：2026年Q1财务报告脱敏",
            key="batch_name_input",
        )
    with batch_cols[1]:
        # 自动生成批次ID，每次 rerun 重新生成
        batch_id = _generate_batch_id()
        st.text_input(
            "批次ID（自动生成）",
            value=batch_id,
            disabled=True,
            key="batch_id_display",
        )

    exec_cols = st.columns(3)
    with exec_cols[0]:
        execute_btn = st.button(
            "🚀 执行脱敏",
            type="primary",
            width="stretch",
        )
    with exec_cols[1]:
        re_detect_btn = st.button(
            "🔄 重新检测",
            width="stretch",
        )

    if re_detect_btn:
        for key in ["detection_results", "file_results", "user_selections", "learn_set"]:
            if key in st.session_state:
                del st.session_state[key]
        st.rerun()

    if execute_btn:
        with st.spinner("正在执行脱敏..."):
            _run_masking(
                uploaded_files, final_selected, all_results,
                mode, ner_enabled, irreversible, learn_words,
                batch_id, batch_name,
            )


def _render_mask_result():
    """展示脱敏结果页面（持久化到 session_state）"""
    result = st.session_state["mask_result"]

    # 成功横幅
    st.markdown(
        '<div class="success-banner">'
        '<h2>🎉 脱敏完成！</h2>'
        f'<p>成功处理 {result["confirmed_count"]} 项敏感信息</p>'
        f'<p>批次ID: {result["batch_id"]} | 批次名称: {result["batch_name"] or "未命名"}</p>'
        '</div>',
        unsafe_allow_html=True,
    )

    # 映射表信息
    mappings = result.get("mappings", [])
    if mappings:
        with st.expander("📋 脱敏映射表", expanded=False):
            mapping_df = pd.DataFrame([
                {
                    "Token": m["token"],
                    "原文": m["original"],
                    "类别": m.get("type_label", ""),
                    "置信度": f"{m.get('confidence', 0):.2f}",
                }
                for m in mappings
            ])
            st.dataframe(mapping_df, width="stretch", hide_index=True)

    # 下载按钮
    st.markdown("#### 📥 下载脱敏文件")

    dl_cols = st.columns(2)
    with dl_cols[0]:
        st.download_button(
            label="📥 下载脱敏文件 (ZIP)",
            data=result["zip_buffer"],
            file_name=f"{result['batch_id']}_masked.zip",
            mime="application/zip",
            type="primary",
            width="stretch",
        )
    with dl_cols[1]:
        st.download_button(
            label="📋 下载映射表 (JSON)",
            data=result["mapping_data"].encode("utf-8"),
            file_name=f"{result['batch_id']}_mapping.json",
            mime="application/json",
            width="stretch",
        )

    # 返回按钮
    if st.button("🔄 返回重新脱敏", width="stretch"):
        del st.session_state["mask_result"]
        for key in ["detection_results", "file_results", "user_selections", "learn_set"]:
            if key in st.session_state:
                del st.session_state[key]
        st.rerun()


# ──────────────────────────────────────────────
# 标签页2：恢复还原
# ──────────────────────────────────────────────

def _render_restore_tab():
    """渲染恢复还原标签页"""

    # 检查是否有恢复结果需要展示
    if "restore_result" in st.session_state:
        restore_result = st.session_state["restore_result"]

        st.markdown(
            '<div class="success-banner">'
            '<h2>🎉 恢复完成！</h2>'
            f'<p>成功恢复 {restore_result["file_count"]} 个文件</p>'
            '</div>',
            unsafe_allow_html=True,
        )

        st.download_button(
            label="📥 下载恢复文件 (ZIP)",
            data=restore_result["zip_buffer"],
            file_name="restored_files.zip",
            mime="application/zip",
            type="primary",
            width="stretch",
        )

        if st.button("🔄 返回", width="stretch"):
            del st.session_state["restore_result"]
            st.rerun()
        return

    st.markdown("#### 🔓 恢复还原")
    st.caption("上传脱敏后的文件和映射表，将敏感信息还原为原始内容")

    st.markdown("---")

    # 选择恢复方式
    restore_method = st.radio(
        "选择恢复方式",
        options=["从历史记录恢复", "手动上传文件恢复"],
        horizontal=True,
    )

    tokens: Optional[dict] = None

    if restore_method == "从历史记录恢复":
        tokens = _render_history_selector()
    else:
        tokens = _render_manual_upload()

    if tokens is None:
        return

    # 显示映射表预览
    st.markdown("#### 📋 映射表预览")
    preview_items = [
        f"- `{token}` → `{original}`"
        for token, original in list(tokens.items())[:20]
    ]
    st.markdown("\n".join(preview_items))
    if len(tokens) > 20:
        st.caption(f"... 共 {len(tokens)} 条映射")

    st.markdown("---")

    # 上传脱敏后的文件
    st.markdown("#### 📁 上传脱敏后的文件")
    masked_files = st.file_uploader(
        "上传需要恢复的脱敏文件",
        type=["docx", "xlsx", "pptx", "pdf", "txt"],
        accept_multiple_files=True,
        label_visibility="collapsed",
    )

    if not masked_files:
        st.info("📤 请上传需要恢复的脱敏文件")
        return

    # 显示已上传文件
    file_cols = st.columns(min(len(masked_files), 4))
    for i, f in enumerate(masked_files):
        with file_cols[i % len(file_cols)]:
            icon = _file_icon(Path(f.name).suffix)
            size_kb = f.size / 1024
            st.markdown(
                f'<div class="file-card">'
                f'<span class="icon">{icon}</span>'
                f'<div><div class="name">{f.name}</div>'
                f'<div class="size">{size_kb:.1f} KB</div></div>'
                f'</div>',
                unsafe_allow_html=True,
            )

    # 执行恢复按钮
    if st.button("🔓 执行恢复", type="primary", width="stretch"):
        with st.spinner("正在恢复文件..."):
            _run_restore(masked_files, tokens)


def _render_history_selector() -> Optional[dict]:
    """从历史记录中选择批次，返回 tokens dict"""
    records = _load_history()

    if not records:
        st.warning("⚠️ 暂无历史记录，请先执行脱敏操作或选择手动上传方式")
        return None

    st.markdown("#### 📚 历史记录")

    # 构建选择列表
    options = []
    for r in reversed(records):  # 最新的在前
        name_part = f" | {r.batch_name}" if r.batch_name else ""
        options.append(
            f"{r.batch_id}{name_part} | {r.created_at[:19]} | "
            f"{r.file_count}个文件 | {r.mask_count}项脱敏"
        )

    selected_idx = st.selectbox(
        "选择批次",
        options=range(len(options)),
        format_func=lambda i: options[i],
    )

    # 获取选中的记录（倒序索引）
    record = records[-(selected_idx + 1)]

    # 解析映射数据
    try:
        mapping_json = json.loads(record.mapping_data)
        # Pipeline.save_mapping 格式: {"tokens": {token: {token, original, ...}}, "metadata": {...}}
        if isinstance(mapping_json, dict) and "tokens" in mapping_json:
            raw_tokens = mapping_json["tokens"]
        elif isinstance(mapping_json, dict):
            raw_tokens = mapping_json
        elif isinstance(mapping_json, list):
            raw_tokens = mapping_json
        else:
            st.error("映射表格式不正确")
            return None

        # 统一转换为 {token_str: original_str} 格式
        tokens = {}
        if isinstance(raw_tokens, dict):
            for k, v in raw_tokens.items():
                if isinstance(v, dict) and "original" in v:
                    tokens[k] = v["original"]
                elif isinstance(v, str):
                    tokens[k] = v
        elif isinstance(raw_tokens, list):
            for item in raw_tokens:
                if isinstance(item, dict) and "token" in item and "original" in item:
                    tokens[item["token"]] = item["original"]
    except (json.JSONDecodeError, TypeError) as e:
        st.error(f"映射表解析失败: {e}")
        return None

    if not tokens:
        st.warning("该批次没有映射数据")
        return None

    # 显示批次详情
    detail_cols = st.columns(4)
    with detail_cols[0]:
        st.metric("批次ID", record.batch_id)
    with detail_cols[1]:
        st.metric("批次名称", record.batch_name or "未命名")
    with detail_cols[2]:
        st.metric("文件数量", record.file_count)
    with detail_cols[3]:
        st.metric("脱敏项数", record.mask_count)

    return tokens


def _render_manual_upload() -> Optional[dict]:
    """手动上传映射表，返回 tokens dict"""
    st.markdown("#### 📋 上传映射表")
    st.caption("请上传脱敏时生成的 mapping.json 文件")

    mapping_file = st.file_uploader(
        "上传映射表 JSON",
        type=["json"],
        label_visibility="collapsed",
    )

    if not mapping_file:
        st.info("📤 请上传映射表 JSON 文件")
        return None

    try:
        content = mapping_file.read().decode("utf-8")
        mapping_json = json.loads(content)

        # Pipeline.save_mapping 格式: {"tokens": {...}, "metadata": {...}}
        if isinstance(mapping_json, dict) and "tokens" in mapping_json:
            raw_tokens = mapping_json["tokens"]
        elif isinstance(mapping_json, dict):
            raw_tokens = mapping_json
        elif isinstance(mapping_json, list):
            raw_tokens = mapping_json
        else:
            st.error("映射表格式不正确")
            return None

        # 统一转换为 {token_str: original_str} 格式
        tokens = {}
        if isinstance(raw_tokens, dict):
            for k, v in raw_tokens.items():
                if isinstance(v, dict) and "original" in v:
                    tokens[k] = v["original"]
                elif isinstance(v, str):
                    tokens[k] = v
        elif isinstance(raw_tokens, list):
            for item in raw_tokens:
                if isinstance(item, dict) and "token" in item and "original" in item:
                    tokens[item["token"]] = item["original"]

        if not tokens:
            st.warning("映射表为空")
            return None

        st.success(f"✅ 成功加载 {len(tokens)} 条映射")
        return tokens

    except (json.JSONDecodeError, UnicodeDecodeError, TypeError) as e:
        st.error(f"映射表解析失败: {e}")
        return None


def _run_restore(masked_files, tokens: dict):
    """执行恢复流程"""
    tmp_dir = Path(tempfile.mkdtemp())
    output_dir = tmp_dir / "restored"
    output_dir.mkdir(parents=True, exist_ok=True)

    # 保存上传文件到临时目录
    saved_paths = []
    for f in masked_files:
        save_path = tmp_dir / f.name
        with open(save_path, "wb") as fp:
            fp.write(f.read())
        saved_paths.append(save_path)

    # 逐文件恢复
    output_files: List[Path] = []
    for file_path in saved_paths:
        suffix = file_path.suffix.lower()
        if suffix not in {".docx", ".xlsx", ".pptx", ".pdf", ".txt"}:
            st.warning(f"跳过不支持的文件类型: {file_path.name}")
            continue

        output_path = output_dir / f"{file_path.stem}_restored{suffix}"
        result = _unmask_file(file_path, output_path, tokens)
        if result:
            output_files.append(result)

    if not output_files:
        st.error("❌ 未能恢复任何文件")
        return

    # 打包为 ZIP
    zip_buffer = BytesIO()
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for fp in output_files:
            zf.write(fp, fp.name)

    zip_bytes = zip_buffer.getvalue()

    # 存入 session_state
    st.session_state["restore_result"] = {
        "zip_buffer": zip_bytes,
        "file_count": len(output_files),
    }

    st.rerun()

    # 展示结果（rerun 后会到达下面的代码）
    # 注意：由于上面已经 rerun，下面的代码不会执行
    # 结果展示在 _render_restore_tab 中检查 restore_result


# ──────────────────────────────────────────────
# 检测流程
# ──────────────────────────────────────────────

def _run_detection(uploaded_files, mode: str, ner_enabled: bool):
    """执行检测流程"""
    # 加载配置
    cfg = _load_config(mode)
    cfg.ner.enabled = ner_enabled

    pipeline = Pipeline(cfg)

    # 保存上传文件到临时目录
    tmp_dir = Path(tempfile.mkdtemp())
    saved_paths = []
    for f in uploaded_files:
        save_path = tmp_dir / f.name
        with open(save_path, "wb") as fp:
            fp.write(f.read())
        saved_paths.append(save_path)

    # 逐文件检测
    all_results: List[DetectionResult] = []
    file_results: Dict[str, List[DetectionResult]] = {}

    for file_path in saved_paths:
        try:
            text = _extract_text(file_path)
            if not text.strip():
                continue
            results = pipeline.detector.detect(text, str(file_path))
            results = pipeline.policy.apply(results)
            all_results.extend(results)
            file_results[file_path.name] = results
        except Exception as e:
            st.error(f"检测 {file_path.name} 时出错: {e}")

    # 跨文件去重
    all_results = _dedup_results(all_results)

    # 存入 session_state
    st.session_state["detection_results"] = all_results
    st.session_state["file_results"] = file_results
    st.session_state["tmp_dir"] = str(tmp_dir)
    st.session_state["saved_paths"] = [str(p) for p in saved_paths]

    st.success(f"✅ 检测完成！共发现 **{len(all_results)}** 项敏感信息")
    st.rerun()


# ──────────────────────────────────────────────
# 脱敏执行流程
# ──────────────────────────────────────────────

def _run_masking(
    uploaded_files,
    selected_indices: List[int],
    all_results: List[DetectionResult],
    mode: str,
    ner_enabled: bool,
    irreversible: bool,
    learn_words: bool,
    batch_id: str,
    batch_name: str,
):
    """执行脱敏流程，结果持久化到 session_state"""
    # 加载配置
    cfg = _load_config(mode)
    cfg.ner.enabled = ner_enabled

    pipeline = Pipeline(cfg)

    if irreversible:
        from mask_tool.core.masker import Masker
        from mask_tool.core.tokenizer import TokenGenerator
        pipeline.masker = Masker(TokenGenerator(), irreversible=True)

    # 获取确认的检测结果
    confirmed_results = [all_results[i] for i in selected_indices]
    # 设置状态为 AUTO_MASK
    for r in confirmed_results:
        r.status = DetectionStatus.AUTO_MASK

    # 处理学习词
    learn_set = st.session_state.get("learn_set", set())
    learned_words: Dict[str, List[str]] = {}
    for i in learn_set:
        if i < len(all_results):
            r = all_results[i]
            cat = r.text_type.value
            if cat not in learned_words:
                learned_words[cat] = []
            if r.text not in learned_words[cat]:
                learned_words[cat].append(r.text)

    # 保存上传文件到临时目录
    tmp_dir = Path(st.session_state.get("tmp_dir", tempfile.mkdtemp()))
    saved_paths = []
    for f in uploaded_files:
        save_path = tmp_dir / f.name
        if not save_path.exists():
            with open(save_path, "wb") as fp:
                fp.write(f.read())
        saved_paths.append(save_path)

    # 输出目录
    output_dir = tmp_dir / "output"
    output_dir.mkdir(parents=True, exist_ok=True)

    # 逐文件脱敏
    output_files: List[Path] = []
    for file_path in saved_paths:
        try:
            # 提取文本
            text = _extract_text(file_path)
            if not text.strip():
                continue

            # 获取该文件相关的确认结果
            file_confirmed = [
                r for r in confirmed_results
                if Path(r.location.file).name == file_path.name
                or not r.location.file  # 无文件信息的也处理
            ]

            if not file_confirmed:
                # 如果没有特定于该文件的结果，使用所有确认结果
                file_confirmed = confirmed_results

            # 执行脱敏
            result_path = _do_mask_file(file_path, output_dir, pipeline, file_confirmed)
            if result_path:
                output_files.append(result_path)

        except Exception as e:
            st.error(f"脱敏 {file_path.name} 时出错: {e}")

    # 保存映射表
    mapping_path = output_dir / "mapping.json"
    pipeline.save_mapping(mapping_path)

    # 学习新词
    if learn_words and learned_words:
        _save_learned_words(learned_words, cfg)

    # ── 生成结果并持久化 ──
    if output_files:
        # 打包为 ZIP
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for fp in output_files:
                zf.write(fp, fp.name)
            # 同时包含映射表
            if mapping_path.exists():
                zf.write(mapping_path, "mapping.json")

        zip_bytes = zip_buffer.getvalue()

        # 读取映射数据
        if mapping_path.exists():
            with open(mapping_path, "r", encoding="utf-8") as f:
                mapping_data = f.read()
        else:
            mapping_data = "{}"

        # 构建映射列表（用于展示）
        mappings = []
        for m in pipeline.masker.get_mappings():
            mappings.append({
                "token": m.token,
                "original": m.original,
                "type_label": TYPE_LABELS.get(m.text_type, m.text_type.value),
                "confidence": m.confidence,
            })

        # 持久化到 session_state
        st.session_state["mask_result"] = {
            "zip_buffer": zip_bytes,
            "mapping_data": mapping_data,
            "mappings": mappings,
            "output_files": [str(p) for p in output_files],
            "confirmed_count": len(confirmed_results),
            "batch_id": batch_id,
            "batch_name": batch_name,
        }

        # 保存批次历史
        history_record = BatchRecord(
            batch_id=batch_id,
            batch_name=batch_name or "",
            created_at=datetime.now().isoformat(),
            file_count=len(output_files),
            mask_count=len(confirmed_results),
            mapping_data=mapping_data,
        )
        _add_history(history_record)

        st.rerun()
    else:
        st.warning("⚠️ 未能生成脱敏文件")


def _save_learned_words(learned: dict, config: MaskConfig):
    """将学习到的词追加到词库文件"""
    lexicon_path = Path(config.lexicon_path)
    if not lexicon_path.exists():
        return

    with open(lexicon_path, "r", encoding="utf-8") as f:
        existing = yaml.safe_load(f) or {}

    new_count = 0
    for category, words in learned.items():
        if category not in existing:
            existing[category] = []
        for word in words:
            if word not in existing[category]:
                existing[category].append(word)
                new_count += 1

    if new_count > 0:
        with open(lexicon_path, "w", encoding="utf-8") as f:
            yaml.dump(existing, f, allow_unicode=True, default_flow_style=False)


# ──────────────────────────────────────────────
# 主应用
# ──────────────────────────────────────────────

def main():
    _inject_css()

    # 侧边栏
    mode, ner_enabled, irreversible, learn_words = render_sidebar()

    # 标题
    st.title("🔒 文件脱敏工具")
    st.caption("上传文件 → 智能检测敏感信息 → 交互式确认 → 一键脱敏下载")

    # 标签页
    tab1, tab2 = st.tabs(["🔒 脱敏处理", "🔓 恢复还原"])

    with tab1:
        _render_masking_tab(mode, ner_enabled, irreversible, learn_words)

    with tab2:
        _render_restore_tab()


# ──────────────────────────────────────────────
# 入口
# ──────────────────────────────────────────────

def run_web():
    """CLI 入口：启动 Streamlit Web 界面

    用法：
        mask-tool-web              # 默认启动
        mask-tool-web --port 8080  # 指定端口
    """
    import sys
    import subprocess

    app_file = Path(__file__).resolve()
    args = [sys.executable, "-m", "streamlit", "run", str(app_file)]
    # 透传额外参数（如 --port, --server.headless 等）
    if len(sys.argv) > 1:
        args.extend(sys.argv[1:])
    subprocess.run(args)


if __name__ == "__main__":
    main()
